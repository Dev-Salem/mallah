from __future__ import annotations

from pathlib import Path
from textwrap import wrap

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUTPUT_DIR = ROOT / "output"
ASSETS_DIR = ROOT / "assets"
OUTPUT_PPTX = OUTPUT_DIR / "Mallah_Award_Winning_Graduation_Deck.pptx"
OUTPUT_BRIEF = OUTPUT_DIR / "Mallah_Award_Winning_Graduation_Deck.md"


BG = RGBColor(17, 20, 24)
PANEL = RGBColor(30, 34, 39)
PANEL_ALT = RGBColor(24, 28, 33)
TEXT = RGBColor(242, 236, 224)
MUTED = RGBColor(176, 170, 158)
ACCENT = RGBColor(224, 130, 58)
ACCENT_SOFT = RGBColor(87, 54, 33)
LINE = RGBColor(62, 69, 77)
SUCCESS = RGBColor(79, 184, 141)
WARNING = RGBColor(216, 171, 87)

TITLE_FONT = "Aptos Display"
BODY_FONT = "Aptos"
MONO_FONT = "Aptos Mono"


REFERENCES = [
    "[1] A. Burke, Digital Technologies in Career Guidance for Youth: Opportunities and Challenges, OECD, 2024.",
    "[2] OECD, Observatory on Digital Technologies in Career Guidance for Youth (ODiCY), 2024.",
    "[3] S. Bankins and S. Formosa, Navigating career stages in the age of artificial intelligence: A systematic interdisciplinary review and agenda for future research, Business Horizons, 2024.",
    "[4] R. Buchanan, Exploring the Role of AI in Career Access Through a Social Justice Lens, Nordic Journal of Transitions, Careers and Guidance, 2025.",
    "[5] Kingdom of Saudi Arabia, Saudi Vision 2030, 2025.",
    "[6] Vision 2030, Human Capability Development Program, 2025.",
    "[7] MCIT, Digital Economy Policy in the Kingdom of Saudi Arabia, 2020.",
    "[8] Digital Government Authority (DGA), Digital Transformation Strategies Across Saudi Arabia, v1.0, 2024.",
]


SLIDES = [
    {
        "number": "01",
        "label": "Graduation Project",
        "title": "Mallah: From Confusion to Job Readiness",
        "support": "An AI-guided platform that helps CS and IT learners choose a path, build proof, and move toward real opportunities with more structure and less wasted effort.",
        "notes": (
            "Open with the core promise. Mallah is not another content site; it is a system that turns uncertainty into a guided progression from learning to employability. "
            "Name the team, supervisor, and university, then move quickly into the problem."
        ),
        "meta_left": "Presented by Abdulaziz F. Alotaibi · Salem S. Hamed · Nader F. Alotaibi · Jamal A. Sweidan",
        "meta_right": "Supervisor: Dr. Mohammed Elfaki · Shaqra University",
        "kind": "title",
    },
    {
        "number": "02",
        "label": "Problem",
        "title": "Graduates do not need more content.",
        "support": "They need direction, proof, and a clear moment when they know they are ready to apply.",
        "cards": [
            ("Theory", "University builds knowledge. Hiring asks for visible execution."),
            ("Noise", "The internet offers thousands of disconnected tutorials and no sequence."),
            ("Readiness", "Learners cannot tell whether to study more or start applying."),
        ],
        "notes": (
            "This slide should feel human, not abstract. Emphasize that the failure is not effort; the failure is fragmentation. "
            "Set up the rest of the deck around the missing link between education and employability."
        ),
        "kind": "cards",
    },
    {
        "number": "03",
        "label": "Saudi Context",
        "title": "This gap matters even more now.",
        "support": "Saudi digital transformation creates demand for talent, but many learners still reach graduation without a structured bridge into entry-level technical work.",
        "stats": [
            ("Local Need", "Digital employability is a national priority, not a side topic."),
            ("Graduate Reality", "Students still face path confusion, weak portfolios, and unclear readiness."),
            ("Product Opportunity", "A guided bilingual platform fits both the market need and the learner context."),
        ],
        "notes": (
            "Tie the problem to national relevance without sounding generic. Use the references as proof that digital guidance and employability matter at a systems level, then bring it back to the learner experience."
        ),
        "kind": "stats",
    },
    {
        "number": "04",
        "label": "Product Definition",
        "title": "Mallah connects the full journey.",
        "support": "One platform links onboarding, roadmap learning, project evidence, resume building, job analysis, and application tracking into a single employability workflow.",
        "flow": [
            "Onboarding",
            "Path Match",
            "Roadmap",
            "Projects",
            "Resume",
            "Job Analysis",
            "Tracking",
        ],
        "notes": (
            "Define Mallah in one sentence and then walk left to right. The message is that the value comes from connection. Each module becomes more useful because it is connected to the others."
        ),
        "kind": "flow",
    },
    {
        "number": "05",
        "label": "Learner Journey",
        "title": "Four stages. One destination.",
        "support": "The learner moves through a simple progression: understand where to start, follow a path, build proof, then approach the market with evidence.",
        "stages": [
            ("01", "Assessment", "Profile, goals, time, confidence"),
            ("02", "Pathing", "One of four structured technical tracks"),
            ("03", "Execution", "Topics, milestones, portfolio, resume"),
            ("04", "Readiness", "Opportunity analysis and application action"),
        ],
        "notes": (
            "This is the bridge slide between problem and product proof. Keep it calm and visual. The audience should be able to repeat the system back after this slide."
        ),
        "kind": "stages",
    },
    {
        "number": "06",
        "label": "Platform Overview",
        "title": "A working system, not a feature list.",
        "support": "These modules are already designed to reinforce each other instead of living as isolated tools.",
        "screens": [
            ("dashboard.png", "Dashboard / command center"),
            ("roadmap.png", "Roadmap / topic viewer"),
            ("portfolio.png", "Portfolio / project proof"),
            ("resume-builder.png", "Resume builder"),
            ("opportunity-analyzer.png", "Opportunity analyzer"),
            ("tracker.png", "Application tracker"),
        ],
        "notes": (
            "Treat this as the visual proof slide. If screenshots are available, pause long enough for the judges to register that the product is real and broad."
        ),
        "kind": "collage",
    },
    {
        "number": "07",
        "label": "Deep Dive 1",
        "title": "The roadmap is the engine.",
        "support": "Mallah organizes each path into stages and topics, tracks progress, and uses project milestones to turn passive study into visible advancement.",
        "bullets": [
            "4 predefined technical paths",
            "Path -> Stage -> Topic structure",
            "Project-gated stage progression",
            "Topic-level progress and AI lesson support",
        ],
        "screen": ("roadmap.png", "Roadmap / Topic Viewer"),
        "notes": (
            "Lead with learner value first: clarity and sequence. Then mention the internal logic: stage structure, topic progress, and milestone gating make the roadmap a real system rather than a static checklist."
        ),
        "kind": "deep_dive",
    },
    {
        "number": "08",
        "label": "Deep Dive 2",
        "title": "Learning becomes recruiter-visible proof.",
        "support": "Portfolio Hub and Resume Builder turn completed work into evidence that can be shown, shared, and tailored for opportunities.",
        "bullets": [
            "Verified skills unlocked from roadmap and projects",
            "Project records with public visibility controls",
            "ATS-oriented resume builder linked to real learner data",
            "A direct bridge from proof to application assets",
        ],
        "screen_pair": [
            ("portfolio.png", "Portfolio / project evidence"),
            ("resume-builder.png", "Resume builder editor"),
        ],
        "notes": (
            "This is where you show the deck that Mallah does not stop at learning. It converts progress into visible proof. That is a major difference between education tooling and employability tooling."
        ),
        "kind": "deep_dive_pair",
    },
    {
        "number": "09",
        "label": "Deep Dive 3",
        "title": "Job readiness becomes measurable.",
        "support": "The Opportunity Analyzer and Application Tracker connect jobs to skills, missing topics, action plans, and follow-through.",
        "bullets": [
            "Curated opportunity feed plus manual job description analysis",
            "Match score, missing skills, and roadmap-linked next steps",
            "Resume and portfolio relevance in one workflow",
            "Application stages tracked from saved to offer",
        ],
        "screen_pair": [
            ("opportunity-analyzer.png", "Opportunity Analyzer results"),
            ("tracker.png", "Application Tracker"),
        ],
        "notes": (
            "This is the strongest differentiator slide. Make the contrast explicit: most tools either help you learn or help you apply. Mallah connects the gap between the two."
        ),
        "kind": "deep_dive_pair",
    },
    {
        "number": "10",
        "label": "Architecture",
        "title": "The system was built for reliability.",
        "support": "The technical architecture favors modularity, secure data handling, and clear boundaries between learner-facing features and platform logic.",
        "architecture_columns": [
            ("Frontend", ["Next.js 16", "React 19", "App Router", "Bilingual UI"]),
            ("Application", ["Feature-based modules", "Server actions", "Zod validation", "Rule-based logic first"]),
            ("Data + AI", ["Supabase / PostgreSQL", "RLS-aware access", "OpenAI for language tasks", "Admin visibility"]),
        ],
        "notes": (
            "Technical judges need to trust the implementation without being buried in stack noise. Mention that AI is used by exception, not everywhere, and that the architecture reflects that choice."
        ),
        "kind": "architecture",
    },
    {
        "number": "11",
        "label": "Methodology",
        "title": "We built it iteratively and deliberately.",
        "support": "The project followed an Agile path from problem framing and analysis through design, implementation, testing, and live deployment.",
        "timeline": [
            "Planning",
            "Requirements",
            "System Design",
            "Implementation",
            "Testing",
            "Deployment",
        ],
        "notes": (
            "Keep this visual and compressed. The goal is to show engineering discipline and project maturity, not to read a process report."
        ),
        "kind": "timeline",
    },
    {
        "number": "12",
        "label": "What Was Built",
        "title": "This is already a complete product.",
        "support": "Mallah is not a single prototype screen. It is a connected platform with breadth across learner workflow, administration, and bilingual delivery.",
        "stats_big": [
            ("4", "career paths"),
            ("8", "learner-facing features"),
            ("2", "languages"),
            ("1", "admin control layer"),
        ],
        "screen": ("dashboard.png", "Dashboard or landing state"),
        "notes": (
            "This slide should make the judges feel scale. Mention the bilingual experience, the admin layer, and the fact that the platform covers the full learning-to-application loop."
        ),
        "kind": "completion",
    },
    {
        "number": "13",
        "label": "Differentiation",
        "title": "Mallah closes the loop others leave open.",
        "support": "Most tools solve one fragment. Mallah links decision-making, skill-building, evidence, job analysis, and tracking in a single learner workflow.",
        "compare": [
            ("Learning Platform", "Teaches content", "Does not prove readiness"),
            ("Resume Scanner", "Optimizes keywords", "Does not guide learning"),
            ("Job Tracker", "Tracks applications", "Does not build missing capability"),
            ("Mallah", "Connects all four", "Turns gaps into action"),
        ],
        "notes": (
            "Be confident here. The goal is not to insult other tools; it is to show that Mallah is valuable because it integrates what learners normally patch together manually."
        ),
        "kind": "compare",
    },
    {
        "number": "14",
        "label": "Closing",
        "title": "From confusion to clearer readiness.",
        "support": "Mallah gives beginners structure, evidence, and a more honest view of what to learn next, what they can already prove, and when they are ready to move.",
        "closing_points": [
            "Clearer starting point",
            "Structured progression",
            "Real project evidence",
            "Better job-facing decisions",
        ],
        "notes": (
            "Return to the human problem from slide two. End on the learner outcome, not the tech stack. This should sound like a completed, defendable system with a meaningful purpose."
        ),
        "kind": "closing",
    },
    {
        "number": "15",
        "label": "Appendix",
        "title": "References",
        "support": "Primary references preserved from the original proposal deck.",
        "references": REFERENCES,
        "notes": "Use this slide only if asked for sources or context during Q&A.",
        "kind": "references",
    },
    {
        "number": "16",
        "label": "Appendix",
        "title": "Questions?",
        "support": "Thank you for listening.",
        "notes": "Pause here and invite questions. Keep the final tone calm and assured.",
        "kind": "qa",
    },
]


def add_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

    left_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.16), Inches(7.5)
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = ACCENT
    left_bar.line.fill.background()

    top_line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.25), Inches(0.35), Inches(12.7), Inches(0.02)
    )
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = LINE
    top_line.line.fill.background()

    for x in range(1, 13):
        line = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(0.5), Inches(0.01), Inches(6.5)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(28, 31, 36)
        line.line.fill.background()


def add_textbox(slide, left, top, width, height, text, *, font_name=BODY_FONT, size=20, color=TEXT,
                bold=False, align=PP_ALIGN.LEFT, mono=False, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = MONO_FONT if mono else font_name
    font.size = Pt(size)
    font.color.rgb = color
    font.bold = bold
    font.italic = italic
    return box


def add_label(slide, label: str, number: str) -> None:
    add_textbox(
        slide,
        Inches(0.55),
        Inches(0.45),
        Inches(2.5),
        Inches(0.35),
        f"{number}  {label.upper()}",
        size=11,
        color=ACCENT,
        bold=True,
        mono=True,
    )


def add_title(slide, title: str, support: str, *, title_top=0.95) -> None:
    add_textbox(
        slide,
        Inches(0.55),
        Inches(title_top),
        Inches(6.1),
        Inches(1.5),
        title,
        font_name=TITLE_FONT,
        size=28,
        color=TEXT,
        bold=True,
    )
    add_textbox(
        slide,
        Inches(0.58),
        Inches(title_top + 1.05),
        Inches(6.0),
        Inches(1.0),
        support,
        size=16,
        color=MUTED,
    )


def add_panel(slide, left, top, width, height, *, fill=PANEL, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    panel = slide.shapes.add_shape(shape_type, left, top, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = fill
    panel.line.color.rgb = LINE
    return panel


def asset_path(name: str) -> Path | None:
    path = ASSETS_DIR / name
    return path if path.exists() else None


def add_image_or_placeholder(slide, left, top, width, height, name: str, label: str) -> None:
    path = asset_path(name)
    if path:
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)
        outline = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
        outline.fill.background()
        outline.line.color.rgb = LINE
        return

    panel = add_panel(slide, left, top, width, height, fill=PANEL_ALT)
    add_textbox(
        slide,
        left + Inches(0.22),
        top + Inches(0.22),
        width - Inches(0.44),
        Inches(0.25),
        "SCREENSHOT PLACEHOLDER",
        size=10,
        color=ACCENT,
        bold=True,
        mono=True,
    )
    add_textbox(
        slide,
        left + Inches(0.22),
        top + Inches(0.70),
        width - Inches(0.44),
        Inches(0.8),
        label,
        font_name=TITLE_FONT,
        size=18,
        color=TEXT,
        bold=True,
    )
    add_textbox(
        slide,
        left + Inches(0.22),
        top + Inches(1.28),
        width - Inches(0.44),
        height - Inches(1.5),
        "Replace this frame with a real product capture. The layout is already sized for a premium UI screenshot.",
        size=13,
        color=MUTED,
    )
    panel.line.color.rgb = LINE


def add_notes(slide, notes: str) -> None:
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.clear()
    notes_frame.text = notes


def render_title_slide(prs: Presentation, data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_label(slide, data["label"], data["number"])

    accent_block = add_panel(slide, Inches(0.58), Inches(1.05), Inches(0.95), Inches(3.8), fill=ACCENT_SOFT)
    accent_block.line.fill.background()
    add_textbox(slide, Inches(0.8), Inches(1.28), Inches(0.5), Inches(3.2), "M\nA\nL\nL\nA\nH", size=20, color=ACCENT, bold=True, mono=True, align=PP_ALIGN.CENTER)

    add_textbox(
        slide,
        Inches(1.8),
        Inches(1.1),
        Inches(7.8),
        Inches(1.4),
        data["title"],
        font_name=TITLE_FONT,
        size=30,
        color=TEXT,
        bold=True,
    )
    add_textbox(slide, Inches(1.84), Inches(2.35), Inches(6.9), Inches(1.2), data["support"], size=16, color=MUTED)
    add_textbox(slide, Inches(1.85), Inches(4.65), Inches(4.9), Inches(0.5), data["meta_left"], size=12, color=TEXT)
    add_textbox(slide, Inches(7.05), Inches(4.65), Inches(5.25), Inches(0.5), data["meta_right"], size=12, color=MUTED)

    quote = add_panel(slide, Inches(7.1), Inches(1.38), Inches(4.75), Inches(2.35), fill=PANEL)
    add_textbox(slide, Inches(7.38), Inches(1.65), Inches(4.1), Inches(0.35), "THESIS POSITION", size=10, color=ACCENT, bold=True, mono=True)
    add_textbox(
        slide,
        Inches(7.38),
        Inches(2.05),
        Inches(4.0),
        Inches(1.4),
        "A product defense deck for a platform that bridges learning, proof, and employability.",
        font_name=TITLE_FONT,
        size=20,
        color=TEXT,
        bold=True,
    )
    quote.line.color.rgb = LINE
    add_notes(slide, data["notes"])


def render_cards_slide(prs: Presentation, data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_label(slide, data["label"], data["number"])
    add_title(slide, data["title"], data["support"])

    card_lefts = [0.6, 4.35, 8.1]
    for idx, (title, body) in enumerate(data["cards"]):
        add_panel(slide, Inches(card_lefts[idx]), Inches(3.05), Inches(3.0), Inches(2.55), fill=PANEL_ALT)
        add_textbox(slide, Inches(card_lefts[idx] + 0.18), Inches(3.3), Inches(0.5), Inches(0.35), f"0{idx + 1}", size=12, color=ACCENT, bold=True, mono=True)
        add_textbox(slide, Inches(card_lefts[idx] + 0.18), Inches(3.72), Inches(2.5), Inches(0.45), title, font_name=TITLE_FONT, size=20, color=TEXT, bold=True)
        add_textbox(slide, Inches(card_lefts[idx] + 0.18), Inches(4.22), Inches(2.55), Inches(1.05), body, size=14, color=MUTED)

    add_notes(slide, data["notes"])


def render_stats_slide(prs: Presentation, data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_label(slide, data["label"], data["number"])
    add_title(slide, data["title"], data["support"])

    for idx, (title, body) in enumerate(data["stats"]):
        top = 3.05 + idx * 1.15
        add_panel(slide, Inches(0.6), Inches(top), Inches(11.7), Inches(0.92), fill=PANEL_ALT, radius=False)
        add_textbox(slide, Inches(0.82), Inches(top + 0.15), Inches(2.5), Inches(0.3), title.upper(), size=11, color=ACCENT, bold=True, mono=True)
        add_textbox(slide, Inches(3.0), Inches(top + 0.12), Inches(8.8), Inches(0.45), body, size=15, color=TEXT)

    add_notes(slide, data["notes"])


def render_flow_slide(prs: Presentation, data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_label(slide, data["label"], data["number"])
    add_title(slide, data["title"], data["support"])

    x = 0.6
    for idx, step in enumerate(data["flow"]):
        box_w = Inches(1.55)
        add_panel(slide, Inches(x), Inches(3.4), box_w, Inches(1.3), fill=PANEL_ALT)
        add_textbox(slide, Inches(x + 0.14), Inches(3.6), Inches(1.25), Inches(0.55), step, font_name=TITLE_FONT, size=18, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
        if idx < len(data["flow"]) - 1:
            add_textbox(slide, Inches(x + 1.6), Inches(3.82), Inches(0.45), Inches(0.3), "→", font_name=TITLE_FONT, size=22, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
        x += 1.75

    add_notes(slide, data["notes"])


def render_stages_slide(prs: Presentation, data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_label(slide, data["label"], data["number"])
    add_title(slide, data["title"], data["support"])

    for idx, (num, title, body) in enumerate(data["stages"]):
        left = 0.6 + idx * 3.0
        add_panel(slide, Inches(left), Inches(3.05), Inches(2.55), Inches(2.6), fill=PANEL_ALT)
        add_textbox(slide, Inches(left + 0.18), Inches(3.25), Inches(0.5), Inches(0.3), num, size=12, color=ACCENT, bold=True, mono=True)
        add_textbox(slide, Inches(left + 0.18), Inches(3.62), Inches(2.0), Inches(0.45), title, font_name=TITLE_FONT, size=20, color=TEXT, bold=True)
        add_textbox(slide, Inches(left + 0.18), Inches(4.2), Inches(2.1), Inches(0.95), body, size=13, color=MUTED)

    add_notes(slide, data["notes"])


def render_collage_slide(prs: Presentation, data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_label(slide, data["label"], data["number"])
    add_title(slide, data["title"], data["support"])

    positions = [
        (0.6, 2.85, 3.7, 1.55),
        (4.55, 2.85, 3.7, 1.55),
        (8.5, 2.85, 3.7, 1.55),
        (0.6, 4.55, 3.7, 1.55),
        (4.55, 4.55, 3.7, 1.55),
        (8.5, 4.55, 3.7, 1.55),
    ]
    for (asset, label), (left, top, width, height) in zip(data["screens"], positions):
        add_image_or_placeholder(slide, Inches(left), Inches(top), Inches(width), Inches(height), asset, label)

    add_notes(slide, data["notes"])


def render_deep_dive_slide(prs: Presentation, data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_label(slide, data["label"], data["number"])
    add_title(slide, data["title"], data["support"])

    add_image_or_placeholder(slide, Inches(6.65), Inches(1.3), Inches(5.1), Inches(4.7), data["screen"][0], data["screen"][1])

    for idx, bullet in enumerate(data["bullets"]):
        top = 3.1 + idx * 0.7
        add_textbox(slide, Inches(0.75), Inches(top), Inches(0.25), Inches(0.25), "◈", size=14, color=ACCENT, bold=True)
        add_textbox(slide, Inches(1.05), Inches(top - 0.02), Inches(5.1), Inches(0.35), bullet, size=16, color=TEXT)

    add_notes(slide, data["notes"])


def render_deep_dive_pair_slide(prs: Presentation, data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_label(slide, data["label"], data["number"])
    add_title(slide, data["title"], data["support"], title_top=0.92)

    add_image_or_placeholder(slide, Inches(7.0), Inches(1.45), Inches(2.45), Inches(4.8), data["screen_pair"][0][0], data["screen_pair"][0][1])
    add_image_or_placeholder(slide, Inches(9.65), Inches(1.45), Inches(2.45), Inches(4.8), data["screen_pair"][1][0], data["screen_pair"][1][1])

    for idx, bullet in enumerate(data["bullets"]):
        top = 3.0 + idx * 0.72
        add_textbox(slide, Inches(0.75), Inches(top), Inches(0.25), Inches(0.25), "◈", size=14, color=ACCENT, bold=True)
        add_textbox(slide, Inches(1.05), Inches(top - 0.02), Inches(5.65), Inches(0.4), bullet, size=15, color=TEXT)

    add_notes(slide, data["notes"])


def render_architecture_slide(prs: Presentation, data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_label(slide, data["label"], data["number"])
    add_title(slide, data["title"], data["support"])

    positions = [0.7, 4.15, 7.6]
    widths = [3.0, 3.0, 4.4]
    for idx, (title, items) in enumerate(data["architecture_columns"]):
        add_panel(slide, Inches(positions[idx]), Inches(3.0), Inches(widths[idx]), Inches(2.7), fill=PANEL_ALT)
        add_textbox(slide, Inches(positions[idx] + 0.18), Inches(3.22), Inches(widths[idx] - 0.35), Inches(0.35), title, font_name=TITLE_FONT, size=20, color=TEXT, bold=True)
        for item_idx, item in enumerate(items):
            add_textbox(slide, Inches(positions[idx] + 0.2), Inches(3.7 + item_idx * 0.42), Inches(widths[idx] - 0.35), Inches(0.3), f"• {item}", size=14, color=MUTED)

    add_notes(slide, data["notes"])


def render_timeline_slide(prs: Presentation, data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_label(slide, data["label"], data["number"])
    add_title(slide, data["title"], data["support"])

    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.85), Inches(4.1), Inches(10.9), Inches(0.04)).fill.solid()
    line_shape = slide.shapes[-1]
    line_shape.fill.fore_color.rgb = LINE
    line_shape.line.fill.background()

    for idx, stage in enumerate(data["timeline"]):
        x = 0.95 + idx * 1.85
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(3.75), Inches(0.42), Inches(0.42))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT
        circle.line.fill.background()
        add_textbox(slide, Inches(x - 0.22), Inches(4.35), Inches(0.95), Inches(0.55), stage, size=13, color=TEXT, bold=True, align=PP_ALIGN.CENTER)

    add_notes(slide, data["notes"])


def render_completion_slide(prs: Presentation, data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_label(slide, data["label"], data["number"])
    add_title(slide, data["title"], data["support"])

    add_image_or_placeholder(slide, Inches(7.15), Inches(1.7), Inches(4.75), Inches(3.95), data["screen"][0], data["screen"][1])
    for idx, (value, label) in enumerate(data["stats_big"]):
        left = 0.7 + (idx % 2) * 2.95
        top = 3.05 + (idx // 2) * 1.45
        add_panel(slide, Inches(left), Inches(top), Inches(2.55), Inches(1.08), fill=PANEL_ALT)
        add_textbox(slide, Inches(left + 0.15), Inches(top + 0.1), Inches(0.9), Inches(0.5), value, font_name=TITLE_FONT, size=26, color=ACCENT, bold=True)
        add_textbox(slide, Inches(left + 1.05), Inches(top + 0.18), Inches(1.2), Inches(0.4), label, size=14, color=TEXT)

    add_notes(slide, data["notes"])


def render_compare_slide(prs: Presentation, data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_label(slide, data["label"], data["number"])
    add_title(slide, data["title"], data["support"], title_top=0.92)

    headers = ["System Type", "What it does well", "What it leaves open"]
    header_lefts = [0.75, 3.35, 7.0]
    header_widths = [2.2, 3.35, 4.8]
    for left, width, header in zip(header_lefts, header_widths, headers):
        add_textbox(slide, Inches(left), Inches(2.55), Inches(width), Inches(0.32), header.upper(), size=10, color=ACCENT, bold=True, mono=True)

    for idx, row in enumerate(data["compare"]):
        top = 3.0 + idx * 0.82
        fill = PANEL if row[0] == "Mallah" else PANEL_ALT
        add_panel(slide, Inches(0.65), Inches(top), Inches(11.6), Inches(0.62), fill=fill, radius=False)
        colors = [TEXT, MUTED, MUTED]
        if row[0] == "Mallah":
            colors = [ACCENT, TEXT, SUCCESS]
        add_textbox(slide, Inches(0.85), Inches(top + 0.12), Inches(2.0), Inches(0.25), row[0], size=15, color=colors[0], bold=True)
        add_textbox(slide, Inches(3.45), Inches(top + 0.12), Inches(3.0), Inches(0.25), row[1], size=14, color=colors[1])
        add_textbox(slide, Inches(7.1), Inches(top + 0.12), Inches(4.8), Inches(0.25), row[2], size=14, color=colors[2])

    add_notes(slide, data["notes"])


def render_closing_slide(prs: Presentation, data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_label(slide, data["label"], data["number"])

    add_textbox(slide, Inches(0.7), Inches(1.2), Inches(6.2), Inches(1.3), data["title"], font_name=TITLE_FONT, size=30, color=TEXT, bold=True)
    add_textbox(slide, Inches(0.73), Inches(2.35), Inches(5.9), Inches(1.1), data["support"], size=16, color=MUTED)

    for idx, point in enumerate(data["closing_points"]):
        top = 3.55 + idx * 0.65
        add_textbox(slide, Inches(0.82), Inches(top), Inches(0.3), Inches(0.25), "◈", size=14, color=ACCENT, bold=True)
        add_textbox(slide, Inches(1.15), Inches(top - 0.02), Inches(4.6), Inches(0.3), point, size=18, color=TEXT, bold=True)

    quote = add_panel(slide, Inches(7.1), Inches(1.55), Inches(4.55), Inches(3.7), fill=PANEL_ALT)
    add_textbox(slide, Inches(7.35), Inches(1.82), Inches(4.0), Inches(0.25), "FINAL POSITION", size=10, color=ACCENT, bold=True, mono=True)
    add_textbox(
        slide,
        Inches(7.35),
        Inches(2.2),
        Inches(3.9),
        Inches(2.2),
        "Mallah helps beginners stop guessing, start building, and approach opportunities with more proof and more clarity.",
        font_name=TITLE_FONT,
        size=22,
        color=TEXT,
        bold=True,
    )
    quote.line.color.rgb = LINE
    add_notes(slide, data["notes"])


def render_references_slide(prs: Presentation, data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_label(slide, data["label"], data["number"])
    add_title(slide, data["title"], data["support"], title_top=0.88)

    top = 2.35
    for idx, ref in enumerate(data["references"]):
        chunked = "\n".join(wrap(ref, width=92))
        box_h = 0.42 + 0.18 * chunked.count("\n")
        add_textbox(slide, Inches(0.72), Inches(top), Inches(11.2), Inches(box_h), chunked, size=11, color=TEXT if idx < 4 else MUTED)
        top += box_h + 0.1

    add_notes(slide, data["notes"])


def render_qa_slide(prs: Presentation, data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_label(slide, data["label"], data["number"])
    add_textbox(slide, Inches(0.72), Inches(1.65), Inches(6.0), Inches(1.0), data["title"], font_name=TITLE_FONT, size=34, color=TEXT, bold=True)
    add_textbox(slide, Inches(0.76), Inches(2.9), Inches(4.5), Inches(0.4), data["support"], size=18, color=MUTED)

    add_panel(slide, Inches(7.2), Inches(1.5), Inches(4.1), Inches(3.4), fill=PANEL_ALT)
    add_textbox(slide, Inches(7.5), Inches(1.85), Inches(3.2), Inches(0.4), "MALLAH", font_name=TITLE_FONT, size=28, color=ACCENT, bold=True)
    add_textbox(slide, Inches(7.55), Inches(2.55), Inches(3.0), Inches(1.2), "Graduation Project\nShaqra University", size=18, color=TEXT, bold=True)
    add_textbox(slide, Inches(7.55), Inches(4.05), Inches(3.0), Inches(0.5), "AI-guided career navigation for technical learners.", size=12, color=MUTED)
    add_notes(slide, data["notes"])


def build_deck() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    renderers = {
        "title": render_title_slide,
        "cards": render_cards_slide,
        "stats": render_stats_slide,
        "flow": render_flow_slide,
        "stages": render_stages_slide,
        "collage": render_collage_slide,
        "deep_dive": render_deep_dive_slide,
        "deep_dive_pair": render_deep_dive_pair_slide,
        "architecture": render_architecture_slide,
        "timeline": render_timeline_slide,
        "completion": render_completion_slide,
        "compare": render_compare_slide,
        "closing": render_closing_slide,
        "references": render_references_slide,
        "qa": render_qa_slide,
    }

    for slide_data in SLIDES:
        renderers[slide_data["kind"]](prs, slide_data)

    prs.save(OUTPUT_PPTX)


def build_brief() -> None:
    screenshot_map = [
        ("06", "dashboard.png, roadmap.png, portfolio.png, resume-builder.png, opportunity-analyzer.png, tracker.png"),
        ("07", "roadmap.png"),
        ("08", "portfolio.png, resume-builder.png"),
        ("09", "opportunity-analyzer.png, tracker.png"),
        ("12", "dashboard.png or landing.png"),
    ]

    lines = [
        "# Mallah Award-Winning Graduation Deck",
        "",
        f"Generated file: `{OUTPUT_PPTX.name}`",
        "",
        "## Slide-by-slide copy and speaking prompts",
        "",
    ]

    for slide in SLIDES:
        lines.append(f"### Slide {slide['number']} — {slide['title']}")
        lines.append(f"- Label: `{slide['label']}`")
        lines.append(f"- Support: {slide['support']}")
        lines.append(f"- Speaker note prompt: {slide['notes']}")
        lines.append("")

    lines.extend(
        [
            "## Screenshot capture list",
            "",
            "Place these files in `presentation/assets/` and regenerate the deck to replace the built-in premium placeholders.",
            "",
        ]
    )

    for slide_no, assets in screenshot_map:
        lines.append(f"- Slide {slide_no}: `{assets}`")

    lines.extend(
        [
            "",
            "## Expected asset names",
            "",
            "- `landing.png`",
            "- `dashboard.png`",
            "- `roadmap.png`",
            "- `portfolio.png`",
            "- `resume-builder.png`",
            "- `opportunity-analyzer.png`",
            "- `tracker.png`",
            "- `admin-panel.png`",
            "",
            "## References preserved in appendix",
            "",
        ]
    )

    for ref in REFERENCES:
        lines.append(f"- {ref}")

    OUTPUT_BRIEF.write_text("\n".join(lines), encoding="utf-8")


if __name__ == "__main__":
    build_deck()
    build_brief()
